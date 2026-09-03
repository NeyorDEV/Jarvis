"""
document_generator_resolver.py — Génération de documents bureautiques à la voix.

Inspiré des skills Anthropic (docx/pptx/xlsx/pdf) : le contenu est produit par
le LLM sous forme de JSON structuré, puis mis en forme par les bibliothèques
Python dédiées à chaque format. Les fichiers générés sont enregistrés dans
un dossier dédié sur le Bureau pour rester faciles à retrouver.
"""

import os
import re
import json
import asyncio
import builtins
import unicodedata
from datetime import datetime

from google.genai import types

DOSSIER_SORTIE = os.path.join(os.environ.get("USERPROFILE", ""), "Desktop", "JARVIS_Documents")


def nettoyer_accent(texte: str) -> str:
    return "".join(c for c in unicodedata.normalize('NFD', texte.lower().strip()) if unicodedata.category(c) != 'Mn')


def _assurer_dossier_sortie():
    os.makedirs(DOSSIER_SORTIE, exist_ok=True)


def _nom_fichier_sur(titre: str, extension: str) -> str:
    base = re.sub(r'[^a-zA-Z0-9_-]+', '_', titre.strip()).strip('_') or "document"
    horodatage = datetime.now().strftime("%Y%m%d_%H%M%S")
    return os.path.join(DOSSIER_SORTIE, f"{base}_{horodatage}.{extension}")


async def _generer_contenu_ia(prompt: str) -> dict | None:
    """Demande au LLM un contenu structuré en JSON. Retourne None si indisponible/invalide."""
    if not hasattr(builtins, "client") or not hasattr(builtins, "CHOSEN_MODEL"):
        return None
    try:
        def call_gemini():
            return builtins.client.models.generate_content(
                model=builtins.CHOSEN_MODEL,
                contents=prompt,
                config=types.GenerateContentConfig(response_mime_type="application/json"),
            )
        response = await asyncio.to_thread(call_gemini)
        return json.loads(response.text)
    except Exception as e:
        print(f"[DOC GENERATOR] Erreur génération contenu IA : {e}")
        return None


# ── DOCX ─────────────────────────────────────────────────────────────────────
def _rendre_docx(data: dict, chemin: str):
    import docx
    d = docx.Document()
    d.add_heading(data.get("titre", "Document"), level=0)
    for section in data.get("sections", []):
        d.add_heading(section.get("titre", ""), level=1)
        for p in section.get("paragraphes", []):
            d.add_paragraph(p)
    d.save(chemin)


# ── PPTX ─────────────────────────────────────────────────────────────────────
def _rendre_pptx(data: dict, chemin: str):
    from pptx import Presentation
    prs = Presentation()

    slide_titre = prs.slides.add_slide(prs.slide_layouts[0])
    slide_titre.shapes.title.text = data.get("titre", "Présentation")
    if len(slide_titre.placeholders) > 1:
        slide_titre.placeholders[1].text = "Généré par JARVIS"

    layout_contenu = prs.slide_layouts[1]
    for s in data.get("slides", []):
        slide = prs.slides.add_slide(layout_contenu)
        slide.shapes.title.text = s.get("titre", "")
        points = s.get("points", [])
        if points:
            body = slide.placeholders[1].text_frame
            body.clear()
            body.text = points[0]
            for point in points[1:]:
                para = body.add_paragraph()
                para.text = point
    prs.save(chemin)


# ── XLSX ─────────────────────────────────────────────────────────────────────
def _rendre_xlsx(data: dict, chemin: str):
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    ws.title = (data.get("titre", "Feuille") or "Feuille")[:31]

    colonnes = data.get("colonnes", [])
    if colonnes:
        ws.append(colonnes)
        for cell in ws[1]:
            cell.font = Font(bold=True)

    for ligne in data.get("lignes", []):
        ws.append(ligne)

    for col in ws.columns:
        largeur = max((len(str(c.value)) for c in col if c.value is not None), default=8)
        ws.column_dimensions[col[0].column_letter].width = min(largeur + 2, 40)

    wb.save(chemin)


# ── PDF (génération) ─────────────────────────────────────────────────────────
def _rendre_pdf(data: dict, chemin: str):
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    doc = SimpleDocTemplate(chemin, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(data.get("titre", "Document"), styles["Title"]), Spacer(1, 14)]
    for section in data.get("sections", []):
        story.append(Paragraph(section.get("titre", ""), styles["Heading2"]))
        story.append(Spacer(1, 6))
        for p in section.get("paragraphes", []):
            story.append(Paragraph(p, styles["BodyText"]))
            story.append(Spacer(1, 6))
    doc.build(story)


# ── PDF (fusion de fichiers existants) ──────────────────────────────────────
_DOSSIERS_CONNUS = {
    "bureau": "Desktop",
    "documents": "Documents",
    "telechargements": "Downloads",
    "téléchargements": "Downloads",
}


def _fusionner_pdfs_dossier(nom_dossier_utilisateur: str) -> tuple[str | None, int]:
    from PyPDF2 import PdfMerger

    dossier_reel = _DOSSIERS_CONNUS.get(nettoyer_accent(nom_dossier_utilisateur), None)
    if not dossier_reel:
        return None, 0
    chemin_dossier = os.path.join(os.environ.get("USERPROFILE", ""), dossier_reel)
    if not os.path.isdir(chemin_dossier):
        return None, 0

    fichiers_pdf = sorted(f for f in os.listdir(chemin_dossier) if f.lower().endswith(".pdf"))
    if len(fichiers_pdf) < 2:
        return None, len(fichiers_pdf)

    _assurer_dossier_sortie()
    horodatage = datetime.now().strftime("%Y%m%d_%H%M%S")
    chemin_sortie = os.path.join(DOSSIER_SORTIE, f"fusion_{horodatage}.pdf")

    merger = PdfMerger()
    try:
        for f in fichiers_pdf:
            merger.append(os.path.join(chemin_dossier, f))
        merger.write(chemin_sortie)
    finally:
        merger.close()

    return chemin_sortie, len(fichiers_pdf)


async def resoudre_generation_document(cmd: str):
    """
    Génère un document Word/PowerPoint/Excel/PDF à partir d'une demande vocale/écrite,
    ou fusionne les PDF d'un dossier connu (Bureau/Documents/Téléchargements).
    Exemples :
    - "Jarvis, rédige-moi un document word sur l'histoire de Rome"
    - "Jarvis, fais-moi une présentation powerpoint sur le changement climatique"
    - "Jarvis, crée un tableau excel de mon budget mensuel"
    - "Jarvis, génère un pdf sur les bases de Python"
    - "Jarvis, fusionne les pdf du bureau"
    """
    t = nettoyer_accent(cmd)
    t = re.sub(r'^(jarvis|jervis|jarvys|jervys|gervis)(,)?\s*', '', t).strip()

    # ── FUSION PDF (pas de génération IA nécessaire) ───────────────────────
    if any(v in t for v in ["fusionne", "fusionner", "combine", "combiner", "assemble", "assembler"]) and "pdf" in t:
        for nom_dossier in _DOSSIERS_CONNUS:
            if nom_dossier in t:
                chemin, nb_trouves = await asyncio.to_thread(_fusionner_pdfs_dossier, nom_dossier)
                if chemin:
                    return f"C'est fait mylane, j'ai fusionné {nb_trouves} fichiers PDF en un seul document, enregistré dans votre dossier JARVIS_Documents."
                if nb_trouves < 2:
                    return f"Je n'ai trouvé que {nb_trouves} fichier(s) PDF dans ce dossier, il en faut au moins deux pour fusionner."
                return "Une erreur est survenue lors de la fusion des PDF."
        return "Précisez le dossier contenant les PDF à fusionner : bureau, documents ou téléchargements."

    # ── DÉTECTION DU TYPE DE DOCUMENT À GÉNÉRER ────────────────────────────
    verbes_action = ["cree", "creer", "genere", "generer", "redige", "rediger", "fais", "fabrique", "ecris", "ecrire"]
    if not any(re.search(rf'\b{v}\b', t) for v in verbes_action):
        return None

    if any(m in t for m in ["powerpoint", "presentation", "diaporama", "slides"]):
        type_doc, extension = "pptx", "pptx"
        marqueurs = ["powerpoint", "presentation", "diaporama", "slides"]
    elif any(m in t for m in ["excel", "tableur", "feuille de calcul", "tableau excel"]):
        type_doc, extension = "xlsx", "xlsx"
        marqueurs = ["excel", "tableur", "feuille de calcul", "tableau"]
    elif "pdf" in t:
        type_doc, extension = "pdf", "pdf"
        marqueurs = ["pdf"]
    elif any(m in t for m in ["word", "document", "compte-rendu", "compte rendu", "rapport"]):
        type_doc, extension = "docx", "docx"
        marqueurs = ["word", "document", "compte-rendu", "compte rendu", "rapport"]
    else:
        return None

    # Extraction du sujet : tout ce qui suit "sur"/"de"/"à propos de", sinon le reste de la phrase
    sujet = None
    match = re.search(r'(?:sur|a propos de|au sujet de|concernant)\s+(.+)', t)
    if match:
        sujet = match.group(1).strip()
    else:
        reste = t
        for v in verbes_action:
            reste = re.sub(rf'\b{v}\b', '', reste)
        for m in marqueurs:
            reste = reste.replace(m, '')
        for mot_vide in ["un", "une", "des", "le", "la", "les", "moi", "mylane", "de", "du"]:
            reste = re.sub(rf'\b{mot_vide}\b', '', reste)
        sujet = reste.strip()

    if not sujet or len(sujet) < 2:
        return f"Sur quel sujet voulez-vous que je génère ce {type_doc} ?"

    print(f"[DOC GENERATOR] Génération d'un {type_doc} sur : '{sujet}'")

    if hasattr(builtins, "parler"):
        try:
            builtins.parler(f"Très bien, je génère votre {type_doc} sur {sujet}, un instant.")
        except Exception:
            pass

    if type_doc == "pptx":
        prompt = (
            f"Génère le contenu d'une présentation PowerPoint sur le sujet : \"{sujet}\". "
            "Réponds STRICTEMENT avec un objet JSON de ce format : "
            '{"titre": "Titre de la présentation", "slides": [{"titre": "Titre du slide", "points": ["point 1", "point 2"]}]}. '
            "Prévois entre 5 et 8 slides, avec 3 à 5 points concis par slide. Ne rajoute aucun texte hors du JSON."
        )
    elif type_doc == "xlsx":
        prompt = (
            f"Génère les données d'un tableau Excel sur le sujet : \"{sujet}\". "
            "Réponds STRICTEMENT avec un objet JSON de ce format : "
            '{"titre": "Titre du tableau", "colonnes": ["Colonne 1", "Colonne 2"], "lignes": [["valeur", "valeur"]]}. '
            "Les lignes doivent être cohérentes avec les colonnes. Ne rajoute aucun texte hors du JSON."
        )
    else:  # docx ou pdf : même structure de contenu
        prompt = (
            f"Rédige le contenu structuré d'un document sur le sujet : \"{sujet}\". "
            "Réponds STRICTEMENT avec un objet JSON de ce format : "
            '{"titre": "Titre du document", "sections": [{"titre": "Titre de section", "paragraphes": ["paragraphe 1", "paragraphe 2"]}]}. '
            "Prévois entre 3 et 6 sections avec un contenu clair et informatif. Ne rajoute aucun texte hors du JSON."
        )

    data = await _generer_contenu_ia(prompt)
    if not data or not data.get("titre"):
        return f"Désolé mylane, je n'ai pas réussi à générer le contenu de ce {type_doc}."

    try:
        _assurer_dossier_sortie()
        chemin_sortie = _nom_fichier_sur(data["titre"], extension)
        rendus = {"docx": _rendre_docx, "pptx": _rendre_pptx, "xlsx": _rendre_xlsx, "pdf": _rendre_pdf}
        await asyncio.to_thread(rendus[type_doc], data, chemin_sortie)
    except Exception as e:
        print(f"[DOC GENERATOR] Erreur rendu {type_doc} : {e}")
        return f"Désolé mylane, une erreur est survenue lors de la création du fichier {type_doc}."

    if hasattr(builtins, "send_web_action"):
        try:
            await builtins.send_web_action(
                "ctx_card", title="DOCUMENT GÉNÉRÉ",
                text=f"{data['titre']} ({extension.upper()})", type="info", icon="📄"
            )
        except Exception:
            pass

    return f"C'est fait mylane, j'ai généré \"{data['titre']}\" et enregistré le fichier dans votre dossier JARVIS_Documents sur le Bureau."


builtins.resoudre_generation_document = resoudre_generation_document
